"""
definitions.py
==============
Modelos Pydantic, VectorizationService y todas las funciones
compartidas por los tres mecanismos de vectorización: semiautomático,
automático y de actualización.
"""

# ──────────────────────────── IMPORTS ────────────────────────────
import re
import io
import hashlib
import unicodedata
import asyncio
import requests
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from settings import settings, get_model_info, get_model_dimensions, uses_prefix, get_ocr_languages

from pydantic import BaseModel

# Procesamiento de documentos
import fitz          # PyMuPDF
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl
import struct
import olefile

# Audio/video e imagenes
from faster_whisper import WhisperModel
import easyocr
import yt_dlp

# Vectorización
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from openai import AzureOpenAI
from qdrant_client.models import (
    Distance, VectorParams, PointStruct,
    Filter, FieldCondition, MatchValue
)

from langchain_text_splitters import RecursiveCharacterTextSplitter


# ═══════════════════════════════════════════════════════════════════
# MODELOS PYDANTIC
# ═══════════════════════════════════════════════════════════════════

class JobStatus(BaseModel):
    job_id: str
    status: str
    message: Optional[str] = None
    result: Optional[Dict[str, Any]] = None


class DocumentInfo(BaseModel):
    """Info de un archivo listado desde Moodle o desde colección."""
    filename: str
    filesize: Optional[int] = None
    fileurl: Optional[str] = None
    course_name: Optional[str] = None
    file_type: str = "other"
    is_duplicate: bool = False
    existing_chunks: int = 0


class ValidationSummary(BaseModel):
    """Resumen de validación estructural de un documento."""
    filename: str
    file_type: str = "other"
    total_pages: int
    total_images: int
    total_tables: int
    duration_seconds: Optional[float] = None
    has_non_textual: bool
    warning_message: Optional[str] = None
    is_duplicate: bool = False
    existing_chunks: int = 0


class SemiUploadRequest(BaseModel):
    collection_name: str


class AutoPreviewRequest(BaseModel):
    collection_name: str
    curid: int


class AutoVectorizeRequest(BaseModel):
    collection_name: str
    curid: int
    selected_filenames: List[str]
    source_type: Optional[str] = None
    file_source_types: Optional[Dict[str, str]] = None


class UpdateListRequest(BaseModel):
    collection_name: str


class UpdateVectorizeRequest(BaseModel):
    collection_name: str
    filename_to_replace: str
    # Modo Moodle:
    curid: Optional[int] = None
    moodle_filename: Optional[str] = None
    course_name: Optional[str] = None


class ConfirmVectorizeRequest(BaseModel):
    job_id: str


class PipelineConfig(BaseModel):
    """Configuración del pipeline RAG (Strategy Pattern) para una
    vectorización contra una colección híbrida (dense+sparse). Se construye
    automáticamente en core.py a partir de la configuración activa
    (strategies.registry.get_all_active()) — el usuario no arma esto
    manualmente. Ausente (None) para colecciones legacy, que siguen usando
    VectorizationService sin cambios."""
    etl_document_strategy: Optional[str] = None
    etl_audio_strategy: Optional[str] = None
    contextual_strategy: Optional[str] = None
    dense_strategy: str
    sparse_strategy: str
    source_type: str = "curso_propio"


# ═══════════════════════════════════════════════════════════════════
# UTILIDADES DE TEXTO
# ═══════════════════════════════════════════════════════════════════

def _generate_hash(text: str) -> str:
    normalized = re.sub(r"\s+", " ", text.strip().lower())
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def _sanitize_filename(name: str) -> str:
    """Sanitizar nombre para el sistema de archivos."""
    safe = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-_."
    result = "".join(c if c in safe else "_" for c in name)
    return result[:100]


def _normalize_filename(name: str) -> str:
    """
    Normalizar filename para búsquedas en Qdrant.
    """
    normalized = unicodedata.normalize("NFD", name.strip())
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized


# ═══════════════════════════════════════════════════════════════════
# CLASIFICACIÓN DE TIPO DE ARCHIVO
# ═══════════════════════════════════════════════════════════════════

_EXTENSION_TYPE_MAP: Dict[str, str] = {
    ".pdf":  "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".ppt":  "ppt",
    ".xlsx": "xlsx",
    ".jpg":  "image", ".jpeg": "image", ".png": "image", ".webp": "image",
    ".mp4":  "video", ".mov": "video", ".webm": "video", ".mkv": "video",
    ".mp3":  "audio", ".wav": "audio", ".m4a": "audio", ".ogg": "audio",
}


def classify_file_type(filename: str) -> str:
    """
    Clasifica un archivo por su extensión en:
    pdf | docx | pptx | xlsx | image | video | audio | other
    """
    ext = Path(filename).suffix.lower()
    return _EXTENSION_TYPE_MAP.get(ext, "other")


# Firmas (magic bytes) para validar descargas por tipo. Tipos no listados
# (ej. "other") no se validan por firma — solo se exige que el archivo no
# esté vacío, ya que puede no tener un magic byte conocido.
_MAGIC_BYTES: Dict[str, List[bytes]] = {
    "pdf":   [b"%PDF-"],
    "docx":  [b"PK\x03\x04"],  # OOXML = ZIP
    "pptx":  [b"PK\x03\x04"],
    "xlsx":  [b"PK\x03\x04"],
    "ppt":   [b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"],  # OLE2 compound file
    "image": [b"\xff\xd8\xff", b"\x89PNG\r\n\x1a\n", b"RIFF"],  # jpeg, png, webp(RIFF)
    "audio": [b"ID3", b"\xff\xfb", b"\xff\xf3", b"RIFF", b"OggS"],
    "video": [b"\x00\x00\x00", b"RIFF"],  # mp4/mov "ftyp" boxes vary; se valida por tamaño no vacío
}


def _validate_file_signature(file_type: str, header: bytes) -> bool:
    """
    Verifica que los primeros bytes del archivo coincidan con alguna firma
    conocida para su tipo. Si el tipo no tiene firmas registradas (o es
    "video", cuyo header varía según el contenedor), se acepta con solo
    que el archivo no esté vacío.
    """
    signatures = _MAGIC_BYTES.get(file_type)
    if not signatures or file_type == "video":
        return len(header) > 0
    return any(header.startswith(sig) for sig in signatures)


async def _extract_from_pdf(file_path: str) -> str:
    """Extrae texto de PDF excluyendo tablas, encabezados y pies de página."""
    doc = fitz.open(file_path)
    full_text = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_height = page.height
                header_bottom = page_height * 0.10
                footer_top = page_height * 0.90

                tables = page.find_tables(table_settings={
                    "vertical_strategy": "lines",
                    "horizontal_strategy": "lines",
                    "intersection_tolerance": 4,
                })
                table_bboxes = [t.bbox for t in tables]
                words = page.extract_words(x_tolerance=1, y_tolerance=1)

                page_words = []
                for word in words:
                    wb = (word["x0"], word["top"], word["x1"], word["bottom"])
                    if word["top"] < header_bottom or word["bottom"] > footer_top:
                        continue
                    in_table = any(
                        not (wb[2] < tb[0] or wb[0] > tb[2] or
                             wb[3] < tb[1] or wb[1] > tb[3])
                        for tb in table_bboxes
                    )
                    if not in_table:
                        page_words.append(word["text"])

                if page_words:
                    full_text.append(" ".join(page_words))

    except Exception as e:
        raise ValueError(f"Error extrayendo PDF: {e}")
    finally:
        doc.close()

    return "\n".join(full_text)


async def _extract_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


async def _extract_from_pptx(file_path: str) -> str:
    """Extrae título, viñetas, tablas y notas del orador de cada slide."""
    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells)
                    if row_text.strip(" |"):
                        parts.append(row_text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notas: {notes}")

        if parts:
            slides_text.append(f"[Diapositiva {i}]\n" + "\n".join(parts))

    return "\n\n".join(slides_text)


def _iter_ole_records(data: bytes, start: int, end: int):
    """
    Recorre el arbol de registros binarios del formato PowerPoint 97-2003
    (MS-PPT: cada registro tiene un header de 8 bytes: ver+instance (2B),
    recType (2B), recLen (4B)). recVer==0xF marca un contenedor (se recorre
    recursivamente); cualquier otro valor es un atomo (hoja).
    """
    offset = start
    while offset + 8 <= end:
        ver_instance, rec_type, rec_len = struct.unpack_from("<HHI", data, offset)
        rec_ver = ver_instance & 0x0F
        content_start = offset + 8
        content_end = content_start + rec_len
        if content_end > end:
            break
        yield rec_type, rec_ver, content_start, content_end
        offset = content_end


def _extract_ppt_legacy_text(data: bytes) -> str:
    """
    Extrae el texto de los atomos TextCharsAtom (0x0FA0, UTF-16LE) y
    TextBytesAtom (0x0FA8, un byte por caracter) del stream binario
    "PowerPoint Document" de un .ppt (97-2003).
    """
    texts = []

    def walk(start: int, end: int):
        for rec_type, rec_ver, c_start, c_end in _iter_ole_records(data, start, end):
            if rec_type == 0x0FA0:
                texts.append(data[c_start:c_end].decode("utf-16-le", errors="ignore"))
            elif rec_type == 0x0FA8:
                texts.append(data[c_start:c_end].decode("cp1252", errors="ignore"))
            elif rec_ver == 0x0F:
                walk(c_start, c_end)

    walk(0, len(data))
    cleaned = [t.replace("\x0b", "\n").strip() for t in texts]
    return "\n\n".join(t for t in cleaned if t)


async def _extract_from_ppt_legacy(file_path: str) -> str:
    ole = olefile.OleFileIO(file_path)
    try:
        if not ole.exists("PowerPoint Document"):
            raise ValueError("El archivo .ppt no contiene un stream 'PowerPoint Document' valido")
        data = ole.openstream("PowerPoint Document").read()
    finally:
        ole.close()
    return _extract_ppt_legacy_text(data)


async def _extract_from_xlsx(file_path: str) -> str:
    """Extrae cada hoja como filas de texto plano separadas por ' | '."""
    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
    sheets_text = []

    for sheet in wb.worksheets:
        rows_text = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip() != ""]
            if values:
                rows_text.append(" | ".join(values))
        if rows_text:
            sheets_text.append(f"[Hoja: {sheet.title}]\n" + "\n".join(rows_text))

    wb.close()
    return "\n\n".join(sheets_text)


# ── OCR (imagenes) ──────────────────────────────────────────────────

_ocr_reader: Optional["easyocr.Reader"] = None


def _get_ocr_reader() -> "easyocr.Reader":
    global _ocr_reader
    if _ocr_reader is None:
        print(f"🔄 Cargando modelo OCR (idiomas: {settings.OCR_LANGUAGES})...")
        _ocr_reader = easyocr.Reader(get_ocr_languages(), gpu=False)
    return _ocr_reader


async def _extract_from_image(file_path: str) -> str:
    loop = asyncio.get_event_loop()
    reader = await loop.run_in_executor(None, _get_ocr_reader)
    results = await loop.run_in_executor(
        None, lambda: reader.readtext(file_path, detail=0, paragraph=True)
    )
    return "\n".join(results)


# ── ASR (audio/video) ────────────────────────────────────────────────

_whisper_model: Optional[WhisperModel] = None


def _get_whisper_model() -> WhisperModel:
    global _whisper_model
    if _whisper_model is None:
        print(
            f"🔄 Cargando modelo Whisper ({settings.ASR_MODEL_SIZE}, "
            f"{settings.ASR_DEVICE}/{settings.ASR_COMPUTE_TYPE})..."
        )
        _whisper_model = WhisperModel(
            settings.ASR_MODEL_SIZE,
            device=settings.ASR_DEVICE,
            compute_type=settings.ASR_COMPUTE_TYPE,
        )
    return _whisper_model


def get_media_duration(file_path: str) -> float:
    """Duración en segundos leyendo solo metadata del contenedor (sin decodificar)."""
    try:
        import av
        with av.open(file_path) as container:
            if container.duration:
                return round(container.duration / 1_000_000, 1)
    except Exception as e:
        print(f"⚠️ No se pudo leer duración de {file_path}: {e}")
    return 0.0


def extract_audio_track(file_path: str) -> Path:
    """
    Extrae la pista de audio de un archivo de video a un .wav mono 16kHz
    temporal, usando PyAV (decodifica y remuxa sin depender del binario
    ffmpeg). Usada por las estrategias ETL de audio que requieren subir un
    archivo de audio puro (APIs externas), a diferencia de faster-whisper
    local que decodifica el video directamente.
    """
    import av

    input_container = av.open(file_path)
    try:
        try:
            input_audio_stream = next(
                s for s in input_container.streams if s.type == "audio"
            )
        except StopIteration:
            raise ValueError(f"El archivo '{file_path}' no tiene pista de audio")

        output_path = settings.TEMP_DIR / f"{Path(file_path).stem}_audio.wav"
        output_container = av.open(str(output_path), mode="w")
        try:
            output_stream = output_container.add_stream("pcm_s16le", rate=16000)
            output_stream.layout = "mono"
            resampler = av.AudioResampler(format="s16", layout="mono", rate=16000)

            for frame in input_container.decode(input_audio_stream):
                for resampled in resampler.resample(frame):
                    for packet in output_stream.encode(resampled):
                        output_container.mux(packet)
            for packet in output_stream.encode(None):
                output_container.mux(packet)
        finally:
            output_container.close()
    finally:
        input_container.close()

    return output_path


async def transcribe_audio_or_video(file_path: str) -> List[Dict[str, Any]]:
    """
    Transcribe audio/video con faster-whisper.
    Retorna una lista de segmentos: [{"start": float, "end": float, "text": str}, ...]
    """
    loop = asyncio.get_event_loop()
    model = await loop.run_in_executor(None, _get_whisper_model)

    def _run():
        segments_iter, _info = model.transcribe(file_path, vad_filter=True)
        return [
            {"start": s.start, "end": s.end, "text": s.text.strip()}
            for s in segments_iter if s.text.strip()
        ]

    return await loop.run_in_executor(None, _run)


async def _extract_text(file_path: str) -> str:
    path = Path(file_path)
    file_type = classify_file_type(path.name)
    if file_type == "pdf":
        return await _extract_from_pdf(file_path)
    elif file_type == "docx":
        return await _extract_from_docx(file_path)
    elif file_type == "pptx":
        return await _extract_from_pptx(file_path)
    elif file_type == "ppt":
        return await _extract_from_ppt_legacy(file_path)
    elif file_type == "xlsx":
        return await _extract_from_xlsx(file_path)
    elif file_type == "image":
        return await _extract_from_image(file_path)
    raise ValueError(f"Formato no soportado: {path.suffix}")


async def _normalize_text(text: str) -> str:
    text = re.sub(r"[^\w\s.,;:()/\-%¿?!áéíóúñÁÉÍÓÚÑ]", "", text)
    text = re.sub(r"\x00", "", text)
    text = re.sub(r"[\x01-\x08\x0b-\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"(\w+)-\s*\n\s*(\w+)", r"\1\2", text)
    return text.strip()


async def extract_normalize_and_hash(file_path: str) -> Dict[str, Any]:
    raw = await _extract_text(file_path)
    normalized = await _normalize_text(raw)
    if len(normalized) < 100:
        raise ValueError("Documento sin contenido de texto suficiente")
    return {
        "normalized_text": normalized,
        "document_hash": _generate_hash(normalized),
    }


async def chunk_text(text: str) -> List[Dict[str, Any]]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", ", ", " ", ""],
    )
    chunks = []
    for i, chunk in enumerate(splitter.split_text(text)):
        urls = re.findall(r"https?://\S+", chunk)
        chunks.append({
            "chunk": i,
            "text": chunk.strip(),
            "urls": urls if urls else None,
        })
    return chunks


def chunk_audio_segments(segments: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    Agrupa segmentos consecutivos de una transcripción (con start/end en
    segundos) hasta acercarse a CHUNK_SIZE caracteres, conservando el
    start_time/end_time de cada chunk resultante — así el asistente puede
    citar "video X, minuto mm:ss" en vez de perder los tiempos al
    re-particionar el texto plano.
    """
    chunks = []
    current_text = []
    current_start = None
    current_end = None

    def _flush():
        if not current_text:
            return
        text = " ".join(current_text).strip()
        if text:
            chunks.append({
                "chunk": len(chunks),
                "text": text,
                "urls": None,
                "start_time": current_start,
                "end_time": current_end,
            })

    for seg in segments:
        candidate = " ".join(current_text + [seg["text"]])
        if current_text and len(candidate) > settings.CHUNK_SIZE:
            _flush()
            current_text = [seg["text"]]
            current_start = seg["start"]
            current_end = seg["end"]
        else:
            if current_start is None:
                current_start = seg["start"]
            current_text.append(seg["text"])
            current_end = seg["end"]

    _flush()
    return chunks


# ═══════════════════════════════════════════════════════════════════
# VALIDACIÓN ESTRUCTURAL
# ═══════════════════════════════════════════════════════════════════

async def validate_document(file_path: str) -> Dict[str, Any]:
    """
    Valida estructura del documento.
    """
    path = Path(file_path)
    file_type = classify_file_type(path.name)
    stats = {"total_pages": 0, "total_images": 0, "total_tables": 0, "duration_seconds": None}

    if path.suffix.lower() == ".pdf":
        try:
            with pdfplumber.open(file_path) as pdf:
                stats["total_pages"] = len(pdf.pages)
                for page in pdf.pages:
                    tables = page.find_tables(table_settings={
                        "vertical_strategy": "lines",
                        "horizontal_strategy": "lines",
                        "intersection_tolerance": 5,
                        "snap_tolerance": 5,
                    })
                    stats["total_tables"] += len(tables)
                    relevant = [
                        img for img in page.images
                        if float(img.get("width", 0)) > 50
                        and float(img.get("height", 0)) > 50
                    ]
                    stats["total_images"] += len(relevant)

            doc = fitz.open(file_path)
            for page in doc:
                for img in page.get_images(full=True):
                    if img[2] > 100 and img[3] > 100:
                        stats["total_images"] += 1
            doc.close()
        except Exception as e:
            raise ValueError(f"Error validando PDF: {e}")

    elif path.suffix.lower() == ".docx":
        try:
            doc = DocxDocument(file_path)
            stats["total_tables"] = len(doc.tables)
            stats["total_images"] = sum(
                1 for rel in doc.part.rels.values() if "image" in rel.target_ref
            )
            total_words = sum(len(p.text.split()) for p in doc.paragraphs)
            stats["total_pages"] = max(1, total_words // 500)
        except Exception as e:
            raise ValueError(f"Error validando DOCX: {e}")

    elif file_type == "pptx":
        try:
            prs = Presentation(file_path)
            slides = list(prs.slides)
            stats["total_pages"] = len(slides)
            stats["total_images"] = sum(
                1 for slide in slides for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            stats["total_tables"] = sum(
                1 for slide in slides for shape in slide.shapes if shape.has_table
            )
        except Exception as e:
            raise ValueError(f"Error validando PPTX: {e}")

    elif file_type == "ppt":
        try:
            text = await _extract_from_ppt_legacy(file_path)
            total_words = len(text.split())
            stats["total_pages"] = max(1, total_words // 150)
        except Exception as e:
            raise ValueError(f"Error validando PPT: {e}")

    elif file_type == "xlsx":
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            stats["total_pages"] = len(wb.worksheets)
            wb.close()
        except Exception as e:
            raise ValueError(f"Error validando XLSX: {e}")

    elif file_type == "image":
        stats["total_pages"] = 1

    elif file_type in ("audio", "video"):
        stats["duration_seconds"] = get_media_duration(file_path)

    else:
        raise ValueError(f"Formato no soportado para vectorización: '{path.suffix}'")

    warning = None
    if stats["total_images"] > 0 or stats["total_tables"] > 0:
        parts = []
        if stats["total_images"]:
            parts.append(f"{stats['total_images']} imagen(es)")
        if stats["total_tables"]:
            parts.append(f"{stats['total_tables']} tabla(s)")
        warning = (
            f"Se detectaron {' y '.join(parts)} en el documento. "
            "Solo se vectorizará el contenido textual."
        )
    elif file_type == "image":
        warning = "Imagen procesada con OCR — la calidad del texto depende de la nitidez de la imagen."
    elif file_type in ("audio", "video"):
        warning = "Contenido transcrito automáticamente con reconocimiento de voz — puede contener errores."

    return {**stats, "file_type": file_type, "warning_message": warning}


# ═══════════════════════════════════════════════════════════════════
# VECTORIZATION SERVICE
# ═══════════════════════════════════════════════════════════════════

class VectorizationService:
    """
    Servicio de vectorización ligado a un modelo de embeddings concreto.
    """

    def __init__(self, model_name: str):
        """
        model_name debe existir en AVAILABLE_EMBEDDING_MODELS (settings.py).
        Lanza ValueError si el modelo no está en el catálogo.
        """
        # Valida y obtiene la configuración del modelo
        info = get_model_info(model_name)

        self._model_name: str   = model_name
        self._vector_size: int  = info["dimensions"]
        self._use_prefix: bool  = info.get("use_prefix", False)
        self._use_task: bool    = info.get("use_task", False)
        self._task_passage: str = info.get("task_passage", "retrieval.passage")
        self._task_query: str   = info.get("task_query",   "retrieval.query")
        self._provider: str     = info.get("provider", "local")

        print(f"🔄 Cargando modelo : {self._model_name}")
        print(f"   Dimensión       : {self._vector_size}")
        print(f"   Proveedor       : {self._provider}")

        if self._provider == "azure":
            self._model        = None
            self._azure_client = AzureOpenAI(
                api_key=settings.AZURE_EMBEDDING_KEY,
                azure_endpoint=settings.AZURE_EMBEDDING_ENDPOINT,
                api_version=settings.AZURE_EMBEDDING_API_VERSION,
            )
            self._azure_deployment = settings.AZURE_EMBEDDING_DEPLOYMENT
            print(f"   Endpoint Azure  : {settings.AZURE_EMBEDDING_ENDPOINT}")
            print(f"   Deployment      : {self._azure_deployment}")
        else:
            self._azure_client     = None
            self._azure_deployment = None
            print(f"   Usa prefijos    : {self._use_prefix}")
            print(f"   Usa task LoRA   : {self._use_task}")
            self._model = SentenceTransformer(
                self._model_name,
                trust_remote_code=self._use_task,
            )

        print(f"🔄 Conectando a Qdrant: {settings.QDRANT_HOST}:{settings.QDRANT_PORT}")
        self._client = QdrantClient(
            url=f"http://{settings.QDRANT_HOST}:{settings.QDRANT_PORT}",
            timeout=60,
        )
        print(f"✅ VectorizationService listo [{self._model_name}]")

    # ── Propiedades ───────────────────────────────────────────────

    @property
    def model_name(self) -> str:
        return self._model_name

    @property
    def vector_size(self) -> int:
        return self._vector_size

    # ── Validación de compatibilidad con colección ─────────────────

    def validate_model_for_collection(self, collection_name: str) -> None:
        """
        Verifica que la dimensión del modelo coincida con la colección Qdrant.
        Lanza ValueError descriptivo si hay incompatibilidad.

        Para colecciones híbridas (vectores nombrados dense+sparse) no
        aplica: la dimensión efectiva la determina la estrategia de
        embedding denso activa (ver strategies/registry.py y
        core._build_pipeline_config), no el modelo fijo de este servicio
        — así que se omite la validación y se deja pasar.
        """
        if not self.collection_exists(collection_name):
            # Colección inexistente: la validación la hará el endpoint
            return

        try:
            info = self._client.get_collection(collection_name)
            vectors_config = info.config.params.vectors
        except Exception as e:
            raise ValueError(f"No se pudo consultar la colección '{collection_name}': {e}")

        if isinstance(vectors_config, dict):
            return

        col_dim = vectors_config.size
        if col_dim != self._vector_size:
            raise ValueError(
                f"⚠️  Incompatibilidad de dimensiones: la colección '{collection_name}' "
                f"fue creada con vectores de {col_dim} dims, pero el modelo "
                f"'{self._model_name}' genera {self._vector_size} dims. "
                f"Elige un modelo con {col_dim} dimensiones para esta colección."
            )

    # ── Colecciones ──────────────────────────────────────────────

    def list_collections(self) -> List[str]:
        resp = self._client.get_collections()
        return [c.name for c in resp.collections]

    def collection_exists(self, name: str) -> bool:
        return name in self.list_collections()

    # ── Documentos en colección ──────────────────────────────────

    def list_documents_in_collection(self, collection_name: str) -> List[Dict[str, Any]]:
        """
        Lista los documentos únicos vectorizados en una colección.
        """
        if not self.collection_exists(collection_name):
            return []

        seen: Dict[str, Dict] = {}
        offset = None

        while True:
            results, next_offset = self._client.scroll(
                collection_name=collection_name,
                limit=200,
                offset=offset,
                with_payload=["filename", "total_chunks", "date", "course_name"],
                with_vectors=False,
            )
            for point in results:
                fn = point.payload.get("filename")
                if fn and fn not in seen:
                    seen[fn] = {
                        "filename": fn,
                        "total_chunks": point.payload.get("total_chunks", 0),
                        "date": point.payload.get("date", ""),
                        "course_name": point.payload.get("course_name", ""),
                    }
            if next_offset is None:
                break
            offset = next_offset

        return list(seen.values())

    # ── Check duplicado por filename ──────────────────────────────

    def check_duplicate(self, filename: str, collection_name: str) -> Dict[str, Any]:
        """Verifica si un filename ya está vectorizado en la colección."""
        if not self.collection_exists(collection_name):
            return {"exists": False, "total_chunks": 0, "document_hash": None}

        normalized_filename = _normalize_filename(filename)

        try:
            results, _ = self._client.scroll(
                collection_name=collection_name,
                limit=1,
                scroll_filter=Filter(must=[
                    FieldCondition(
                        key="filename_normalized",
                        match=MatchValue(value=normalized_filename)
                    )
                ]),
                with_payload=["document_hash", "total_chunks"],
                with_vectors=False,
            )

            if not results:
                results, _ = self._client.scroll(
                    collection_name=collection_name,
                    limit=1,
                    scroll_filter=Filter(must=[
                        FieldCondition(key="filename", match=MatchValue(value=filename))
                    ]),
                    with_payload=["document_hash", "total_chunks"],
                    with_vectors=False,
                )

            if not results:
                return {"exists": False, "total_chunks": 0, "document_hash": None}

            doc_hash = results[0].payload.get("document_hash")

            all_results, _ = self._client.scroll(
                collection_name=collection_name,
                limit=10000,
                scroll_filter=Filter(must=[
                    FieldCondition(key="document_hash", match=MatchValue(value=doc_hash))
                ]),
                with_payload=False,
                with_vectors=False,
            )
            return {
                "exists": True,
                "total_chunks": len(all_results),
                "document_hash": doc_hash,
            }
        except Exception as e:
            print(f"⚠️ Error check_duplicate: {e}")
            return {"exists": False, "total_chunks": 0, "document_hash": None}

    # ── Eliminar documento por filename ──────────────────────────

    def delete_document(self, filename: str, collection_name: str) -> Dict[str, Any]:
        check = self.check_duplicate(filename, collection_name)
        if not check["exists"]:
            return {"success": False, "message": f"'{filename}' no existe en la colección"}

        normalized_filename = _normalize_filename(filename)
        doc_hash = check["document_hash"]

        results, _ = self._client.scroll(
            collection_name=collection_name,
            limit=1,
            scroll_filter=Filter(must=[
                FieldCondition(
                    key="filename_normalized",
                    match=MatchValue(value=normalized_filename)
                )
            ]),
            with_payload=["document_hash"],
            with_vectors=False,
        )

        if not results:
            return {"success": False, "message": f"'{filename}' no encontrado en la colección"}

        self._client.delete(
            collection_name=collection_name,
            points_selector=Filter(must=[
                FieldCondition(key="document_hash", match=MatchValue(value=doc_hash))
            ]),
        )
        print(f"🗑️ Eliminado: {filename} ({check['total_chunks']} chunks)")
        return {
            "success": True,
            "deleted_chunks": check["total_chunks"],
            "filename": filename,
        }

    # ── Embed ─────────────────────────────────────────────────────

    async def embed(self, texts: List[str], is_query: bool = False) -> List[List[float]]:
        """
        Genera embeddings para la lista de textos.
        """
        loop = asyncio.get_event_loop()

        if self._provider == "azure":
            # Azure OpenAI Embeddings — batches de 100 (límite recomendado)
            async def _azure_embed() -> List[List[float]]:
                results = []
                batch_size = 100
                for i in range(0, len(texts), batch_size):
                    batch = texts[i:i + batch_size]
                    response = await loop.run_in_executor(
                        None,
                        lambda b=batch: self._azure_client.embeddings.create(
                            model=self._azure_deployment,
                            input=b,
                        )
                    )
                    results.extend([item.embedding for item in response.data])
                return results
            return await _azure_embed()

        if self._use_task:
            # Jina Embeddings v3
            task = self._task_query if is_query else self._task_passage
            return await loop.run_in_executor(
                None,
                lambda: self._model.encode(
                    texts,
                    task=task,
                    normalize_embeddings=True,
                ).tolist(),
            )

        if self._use_prefix:
            # Familia E5
            prefix = "query: " if is_query else "passage: "
            inputs = [prefix + t for t in texts]
        else:
            inputs = texts

        return await loop.run_in_executor(
            None,
            lambda: self._model.encode(inputs, normalize_embeddings=True).tolist(),
        )

    # ── Store ─────────────────────────────────────────────────────

    async def store_vectors(
        self,
        *,
        collection_name: str,
        chunks: List[Dict[str, Any]],
        embeddings: List[List[float]],
        metadata: Dict[str, Any],
    ) -> Dict[str, Any]:
        """
        Almacena vectores en Qdrant.
        metadata debe incluir: document_hash, filename, format,
        total_pages, total_chunks, course_name.

        Valida compatibilidad de dimensiones antes de insertar.
        """
        if not self.collection_exists(collection_name):
            raise ValueError(f"La colección '{collection_name}' no existe en Qdrant.")

        self.validate_model_for_collection(collection_name)

        doc_hash = metadata["document_hash"]
        now = datetime.utcnow().isoformat()
        filename_normalized = _normalize_filename(metadata["filename"])
        points = []

        for chunk, vector in zip(chunks, embeddings):
            pid = hashlib.md5(
                f"{doc_hash}_{chunk['chunk']}".encode()
            ).hexdigest()
            points.append(PointStruct(
                id=pid,
                vector=vector,
                payload={
                    "document_hash":      doc_hash,
                    "filename":           metadata["filename"],
                    "filename_normalized": filename_normalized,
                    "format":             metadata["format"],
                    "file_type":          metadata.get("file_type", metadata["format"]),
                    "total_pages":        metadata["total_pages"],
                    "total_chunks":       metadata["total_chunks"],
                    "course_name":        metadata.get("course_name", "modo manual"),
                    "date":               now,
                    "chunk":              chunk["chunk"],
                    "urls":               chunk["urls"],
                    "text":               chunk["text"],
                    "start_time":         chunk.get("start_time"),
                    "end_time":           chunk.get("end_time"),
                    "embedding_model":    self._model_name,
                },
            ))

        batch_size = 100
        for i in range(0, len(points), batch_size):
            self._client.upsert(
                collection_name=collection_name,
                points=points[i:i + batch_size],
                wait=True,
            )

        print(f"💾 {len(points)} vectores → {collection_name} [{self._model_name}]")
        return {
            "success": True,
            "vectors_stored": len(points),
            "collection": collection_name,
            "document_hash": doc_hash,
            "embedding_model": self._model_name,
        }


# ═══════════════════════════════════════════════════════════════════
# REGISTRY DE SERVICIOS (una instancia por modelo)
# ═══════════════════════════════════════════════════════════════════

_service_registry: Dict[str, VectorizationService] = {}


def get_service(model_name: Optional[str] = None) -> VectorizationService:
    """
    Retorna (o crea) la instancia de VectorizationService para el modelo dado.
    """
    key = model_name or settings.EMBEDDING_MODEL

    # Valida que el modelo exista en el catálogo antes de instanciar
    get_model_info(key)  # lanza ValueError si no está

    if key not in _service_registry:
        print(f"🆕 Instanciando servicio para modelo: {key}")
        _service_registry[key] = VectorizationService(model_name=key)

    return _service_registry[key]


# ═══════════════════════════════════════════════════════════════════
# MOODLE — funciones compartidas
# ═══════════════════════════════════════════════════════════════════

def _add_moodle_token(url: str) -> str:
    """
    Agrega el token de Moodle a una URL de pluginfile.php si no lo trae ya.
    Necesario para que el enlace "Ver" del listado (que el frontend abre
    directo, sin pasar por el backend) pueda abrir el archivo sin sesión.
    """
    if not url or "/pluginfile.php" not in url or "token=" in url:
        return url
    sep = "&" if "?" in url else "?"
    return f"{url}{sep}token={settings.MOODLE_TOKEN}"


_YOUTUBE_RE = re.compile(
    r"(?:youtube\.com/(?:watch\?v=|embed/|shorts/)|youtu\.be/)[\w-]{6,}",
    re.IGNORECASE,
)


def get_moodle_user_fullname(userid: int) -> Optional[str]:
    """Nombre completo del usuario de Moodle (para el saludo del asistente
    público). Usa core_user_get_users_by_field — verificado en vivo que el
    MOODLE_TOKEN actual tiene permiso para esta función. Retorna None si el
    usuario no existe o la llamada falla (el saludo genérico es el
    fallback, no debe romper la ruta pública)."""
    try:
        url = f"{settings.MOODLE_URL}/webservice/rest/server.php"
        params = {
            "wstoken": settings.MOODLE_TOKEN,
            "wsfunction": "core_user_get_users_by_field",
            "moodlewsrestformat": "json",
            "field": "id",
            "values[0]": userid,
        }
        r = requests.get(url, params=params, verify=False, timeout=10)
        r.raise_for_status()
        data = r.json()
        if isinstance(data, dict) and "exception" in data:
            print(f"⚠️ Moodle no pudo resolver el usuario {userid}: {data.get('message')}")
            return None
        if data and isinstance(data, list):
            return data[0].get("fullname")
    except Exception as e:
        print(f"⚠️ Error consultando nombre de usuario Moodle {userid}: {e}")
    return None


def get_moodle_course_name(courseid: int) -> Optional[str]:
    """Nombre completo (`fullname`) del curso de Moodle, para confirmar en
    el paso "Curso" del wizard de asistentes. Usa core_course_get_courses
    (options[ids][]) — verificado en vivo que el MOODLE_TOKEN actual tiene
    permiso para esta función; core_course_get_courses_by_field está
    bloqueada para este token. Retorna None si el curso no existe o la
    llamada falla."""
    try:
        url = f"{settings.MOODLE_URL}/webservice/rest/server.php"
        params = {
            "wstoken": settings.MOODLE_TOKEN,
            "wsfunction": "core_course_get_courses",
            "moodlewsrestformat": "json",
            "options[ids][0]": courseid,
        }
        r = requests.get(url, params=params, verify=False, timeout=10)
        r.raise_for_status()
        data = r.json()
        if isinstance(data, dict) and "exception" in data:
            print(f"⚠️ Moodle no pudo resolver el curso {courseid}: {data.get('message')}")
            return None
        if data and isinstance(data, list):
            return data[0].get("fullname")
    except Exception as e:
        print(f"⚠️ Error consultando nombre de curso Moodle {courseid}: {e}")
    return None


def get_course_resources(curid: int) -> List[Dict[str, Any]]:
    """
    Consulta la API de Moodle y retorna los recursos (módulos) del curso.

    Usa core_course_get_contents en vez de mod_resource_get_resources_by_courses
    porque este último solo devuelve módulos de tipo "Archivo" (mod_resource) y
    omite los archivos dentro de módulos "Carpeta" (mod_folder). core_course_get_contents
    trae todos los tipos de módulo con sus archivos, así que se filtran aquí
    "resource" y "folder" y se normaliza la salida al formato
    {"id", "name", "contentfiles": [...]} que ya consume el resto del sistema.

    También procesa módulos "url" (enlaces externos): si el enlace apunta a
    YouTube se agrega como un archivo de video sintético (source="youtube")
    que el resto del pipeline descarga y transcribe igual que cualquier otro
    video. Enlaces no-YouTube dentro de "url" (formularios, Zoom, biblioteca,
    etc.) se ignoran porque no son contenido descargable/vectorizable.
    """
    url = f"{settings.MOODLE_URL}/webservice/rest/server.php"
    params = {
        "wstoken": settings.MOODLE_TOKEN,
        "wsfunction": "core_course_get_contents",
        "moodlewsrestformat": "json",
        "courseid": curid,
    }
    r = requests.get(url, params=params, verify=False)
    r.raise_for_status()
    data = r.json()

    if isinstance(data, dict) and "exception" in data:
        raise RuntimeError(f"Moodle API error: {data.get('message', data)}")

    resources = []
    for section in data:
        for module in section.get("modules", []):
            modname = module.get("modname")
            if modname not in ("resource", "folder", "url"):
                continue

            if modname == "url":
                contentfiles = []
                for item in module.get("contents", []):
                    link = item.get("fileurl") or ""
                    if item.get("type") == "url" and _YOUTUBE_RE.search(link):
                        module_name = module.get("name") or "video_youtube"
                        contentfiles.append({
                            "filename": f"{module_name}.mp4",
                            "filesize": None,
                            "fileurl":  link,
                            "mimetype": "video/youtube",
                            "source":   "youtube",
                        })
            else:
                contentfiles = [
                    {
                        "filename": f["filename"],
                        "filesize": f.get("filesize"),
                        "fileurl":  _add_moodle_token(f.get("fileurl") or ""),
                        "mimetype": f.get("mimetype"),
                    }
                    for f in module.get("contents", [])
                    if f.get("type") == "file"
                ]

            if not contentfiles:
                continue

            resources.append({
                "id":            module.get("id"),
                "name":          module.get("name"),
                "contentfiles":  contentfiles,
            })

    return resources


def download_pdf_to_temp(fileurl: str, filename: str) -> Path:
    """
    Descarga un archivo (de cualquier tipo soportado) directamente desde el
    webservice de Moodle usando el token (wstoken habilitado para descarga
    de archivos) — ya no requiere login por sesión con usuario/contraseña.
    Valida la firma binaria del archivo descargado según su tipo.
    """
    download_url = fileurl
    if "token=" not in download_url:
        sep = "&" if "?" in download_url else "?"
        download_url = f"{download_url}{sep}token={settings.MOODLE_TOKEN}"

    safe_name = _sanitize_filename(filename)
    dest = settings.TEMP_DIR / safe_name
    file_type = classify_file_type(filename)

    with requests.get(download_url, stream=True, allow_redirects=True, verify=False) as r:
        r.raise_for_status()
        with open(dest, "wb") as f:
            for chunk in r.iter_content(256 * 1024):
                if chunk:
                    f.write(chunk)

    with open(dest, "rb") as f:
        header = f.read(16)
    if not header or not _validate_file_signature(file_type, header):
        dest.unlink(missing_ok=True)
        raise ValueError(f"Archivo corrupto o formato inesperado: {filename}")

    print(f"⬇️ Descargado en temp: {dest}")
    return dest


def download_youtube_audio(youtube_url: str, filename: str) -> Path:
    """
    Descarga solo el audio (mejor calidad disponible) de un video de YouTube
    a la carpeta temporal, usando yt-dlp. No requiere ffmpeg porque no se
    mezcla audio+video (el contenedor descargado se transcribe tal cual con
    faster-whisper, que decodifica cualquier formato vía PyAV).
    """
    safe_stem = Path(_sanitize_filename(filename)).stem
    outtmpl = str(settings.TEMP_DIR / f"{safe_stem}.%(ext)s")

    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": outtmpl,
        "quiet": True,
        "no_warnings": True,
        "noplaylist": True,
    }

    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info = ydl.extract_info(youtube_url, download=True)
        dest = Path(ydl.prepare_filename(info))

    if not dest.exists() or dest.stat().st_size == 0:
        raise ValueError(f"No se pudo descargar el audio de YouTube: {youtube_url}")

    print(f"⬇️ Audio de YouTube descargado: {dest}")
    return dest


def cleanup_temp_file(path: Path) -> None:
    """Elimina un archivo de la carpeta temporal."""
    try:
        path.unlink(missing_ok=True)
        print(f"🗑️ Temp eliminado: {path.name}")
    except Exception as e:
        print(f"⚠️ No se pudo eliminar temp {path}: {e}")


# ═══════════════════════════════════════════════════════════════════
# PIPELINE HÍBRIDO (Strategy Pattern) — colecciones con vectores nombrados
# ═══════════════════════════════════════════════════════════════════

class HybridVectorizationService:
    """Vectorización dirigida por configuración (Strategy Pattern) para
    colecciones híbridas (dense+sparse nombrados). No reemplaza a
    VectorizationService, que sigue sirviendo intacta a las colecciones
    legacy (vector único sin nombre) — ver vectorize_file_with_pages()."""

    def __init__(self, config: "PipelineConfig"):
        self._config = config

    async def process_and_store(
        self,
        *,
        file_path: Path,
        collection_name: str,
        course_name: str,
        total_pages: int,
        original_filename: str = "",
    ) -> Dict[str, Any]:
        from strategies import registry as strategy_registry
        from qdrant_admin import get_qdrant_admin

        stored_filename = original_filename if original_filename else file_path.name
        file_type = classify_file_type(stored_filename)

        etl_document = strategy_registry.get_etl_document_strategy(
            self._config.etl_document_strategy
        )
        etl_audio = strategy_registry.get_etl_audio_strategy(self._config.etl_audio_strategy)
        contextual = strategy_registry.get_contextual_strategy(self._config.contextual_strategy)
        dense = strategy_registry.get_dense_strategy(self._config.dense_strategy)
        sparse = strategy_registry.get_sparse_strategy(self._config.sparse_strategy)

        use_contextual = self._config.contextual_strategy not in (None, "none")

        if file_type in ("audio", "video"):
            segments = await etl_audio.transcribe(str(file_path))
            full_text = " ".join(s["text"] for s in segments)
            if len(full_text.strip()) < 20:
                raise ValueError("Transcripción sin contenido de texto suficiente")
            document_hash = _generate_hash(full_text)
            chunks = chunk_audio_segments(segments)
        else:
            raw_text = await etl_document.extract(str(file_path))
            full_text = await _normalize_text(raw_text)
            if len(full_text) < 100:
                raise ValueError(
                    "El documento no contiene suficiente texto extraíble "
                    "(mínimo 100 caracteres)"
                )
            document_hash = _generate_hash(full_text)
            chunks = await chunk_text(full_text)

        # Contextual Retrieval: cada chunk se enriquece con una cabecera de
        # 50-100 palabras (si la estrategia activa no es "none") y el texto
        # que se vectoriza es cabecera+chunk, aunque el payload conserva el
        # texto crudo por separado.
        embed_texts: List[str] = []
        for chunk in chunks:
            header = ""
            if use_contextual:
                header = await contextual.enrich(full_text, chunk["text"])
            chunk["contextual_header"] = header
            embed_texts.append(f"{header}\n\n{chunk['text']}" if header else chunk["text"])

        dense_vectors = await dense.embed(embed_texts)
        sparse_vectors = await sparse.embed(embed_texts)

        filename_normalized = _normalize_filename(stored_filename)
        now = datetime.utcnow().isoformat()
        points = []
        for chunk, dense_vec, sparse_vec in zip(chunks, dense_vectors, sparse_vectors):
            pid = hashlib.md5(f"{document_hash}_{chunk['chunk']}".encode()).hexdigest()
            points.append({
                "id": pid,
                "dense_vector": dense_vec,
                "sparse_indices": sparse_vec["indices"],
                "sparse_values": sparse_vec["values"],
                "payload": {
                    "document_hash": document_hash,
                    "filename": stored_filename,
                    "filename_normalized": filename_normalized,
                    "format": file_path.suffix.lstrip("."),
                    "file_type": file_type,
                    "total_pages": total_pages,
                    "total_chunks": len(chunks),
                    "course_name": course_name,
                    "date": now,
                    "chunk": chunk["chunk"],
                    "urls": chunk.get("urls"),
                    "text": chunk["text"],
                    "contextual_header": chunk.get("contextual_header") or None,
                    "start_time": chunk.get("start_time"),
                    "end_time": chunk.get("end_time"),
                    "slide_number": chunk.get("slide_number"),
                    "source_url": chunk.get("source_url"),
                    "source_type": self._config.source_type,
                    "embedding_model": self._config.dense_strategy,
                    "sparse_model": self._config.sparse_strategy,
                },
            })

        admin = get_qdrant_admin()
        result = admin.upsert_hybrid_points(collection_name=collection_name, points=points)

        return {
            "success": True,
            "vectors_stored": result["points_upserted"],
            "collection": collection_name,
            "document_hash": document_hash,
            "embedding_model": self._config.dense_strategy,
            "sparse_model": self._config.sparse_strategy,
            "pipeline": "hybrid",
        }


# ═══════════════════════════════════════════════════════════════════
# PIPELINE DE VECTORIZACIÓN (función central compartida)
# ═══════════════════════════════════════════════════════════════════

async def vectorize_file(
    file_path: Path,
    collection_name: str,
    course_name: str,
    service: VectorizationService,
) -> Dict[str, Any]:
    """
    Pipeline completo para un archivo ya guardado en disco:
    extract → normalize → hash → chunk → embed → store.
    """
    extracted = await extract_normalize_and_hash(str(file_path))
    normalized_text = extracted["normalized_text"]
    document_hash   = extracted["document_hash"]

    chunks     = await chunk_text(normalized_text)
    embeddings = await service.embed([c["text"] for c in chunks])

    metadata = {
        "document_hash": document_hash,
        "filename":      file_path.name,
        "format":        file_path.suffix.lstrip("."),
        "total_pages":   0,
        "total_chunks":  len(chunks),
        "course_name":   course_name,
    }

    return await service.store_vectors(
        collection_name=collection_name,
        chunks=chunks,
        embeddings=embeddings,
        metadata=metadata,
    )


async def vectorize_file_with_pages(
    file_path: Path,
    collection_name: str,
    course_name: str,
    total_pages: int,
    service: VectorizationService,
    original_filename: str = "",
    pipeline_config: Optional["PipelineConfig"] = None,
) -> Dict[str, Any]:
    """
    Pipeline completo: extract → normalize → hash → chunk → embed → store.
    Para audio/video usa transcripción + chunking por segmentos (conserva
    timestamps); para el resto usa el pipeline de texto plano habitual.

    Si `pipeline_config` no es None (colección híbrida con la configuración
    RAG activa — ver core.py), delega en HybridVectorizationService y usa
    ETL/contextual/embedding denso+disperso según lo configurado. Si es
    None (colección legacy), el resto de esta función queda intacto: mismo
    comportamiento que antes de introducir el pipeline configurable.
    """
    if pipeline_config is not None:
        hybrid_service = HybridVectorizationService(pipeline_config)
        return await hybrid_service.process_and_store(
            file_path=file_path,
            collection_name=collection_name,
            course_name=course_name,
            total_pages=total_pages,
            original_filename=original_filename,
        )

    stored_filename = original_filename if original_filename else file_path.name
    file_type = classify_file_type(stored_filename)

    if file_type in ("audio", "video"):
        segments = await transcribe_audio_or_video(str(file_path))
        full_text = " ".join(s["text"] for s in segments)
        if len(full_text.strip()) < 20:
            raise ValueError("Transcripción sin contenido de texto suficiente")
        document_hash = _generate_hash(full_text)
        chunks = chunk_audio_segments(segments)
    else:
        extracted = await extract_normalize_and_hash(str(file_path))
        document_hash = extracted["document_hash"]
        chunks = await chunk_text(extracted["normalized_text"])

    embeddings = await service.embed([c["text"] for c in chunks])

    metadata = {
        "document_hash": document_hash,
        "filename":      stored_filename,
        "format":        file_path.suffix.lstrip("."),
        "file_type":     file_type,
        "total_pages":   total_pages,
        "total_chunks":  len(chunks),
        "course_name":   course_name,
    }

    return await service.store_vectors(
        collection_name=collection_name,
        chunks=chunks,
        embeddings=embeddings,
        metadata=metadata,
    )
